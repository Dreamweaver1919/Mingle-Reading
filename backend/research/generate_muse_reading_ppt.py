from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(246, 248, 252)
NAVY = RGBColor(24, 44, 97)
BLUE = RGBColor(60, 110, 196)
TEAL = RGBColor(44, 146, 148)
ORANGE = RGBColor(226, 122, 63)
TEXT = RGBColor(40, 48, 63)
MUTED = RGBColor(92, 104, 128)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(230, 236, 247)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None, dark=False):
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.8), Inches(0.8))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE if dark else NAVY
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(8.5), Inches(0.45))
        sbf = sb.text_frame
        sp = sbf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Microsoft YaHei"
        sr.font.size = Pt(10.5)
        sr.font.color.rgb = LIGHT if dark else MUTED


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Inches(12.15), Inches(7.0), Inches(0.6), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(idx)
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def add_bullets(slide, items, left, top, width, height, font_size=18, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.line_spacing = 1.2
        p.space_after = Pt(7)
    return box


def add_card(slide, left, top, width, height, title, lines, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT
    shape.line.width = Pt(1)
    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.3), Inches(0.35))
    tp = title_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.name = "Microsoft YaHei"
    tr.font.size = Pt(14)
    tr.font.bold = True
    tr.font.color.rgb = accent
    add_bullets(slide, lines, left + Inches(0.18), top + Inches(0.48), width - Inches(0.35), height - Inches(0.55), 11.5)


def add_flow_boxes(slide, labels, left, top, width, height, accent=BLUE):
    gap = Inches(0.15)
    box_w = (width - gap * (len(labels) - 1)) / len(labels)
    for i, label in enumerate(labels):
        x = left + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, box_w, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = accent
        shape.line.width = Pt(1.2)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = TEXT


def add_table_like(slide, headers, rows, left, top, col_widths, row_h=0.45):
    x = left
    y = top
    for i, h in enumerate(headers):
        w = Inches(col_widths[i])
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(row_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = NAVY
        rect.line.color.rgb = WHITE
        tf = rect.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE
        x += w
    cur_y = y + Inches(row_h)
    for ridx, row in enumerate(rows):
        x = left
        for cidx, cell in enumerate(row):
            w = Inches(col_widths[cidx])
            rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, cur_y, w, Inches(row_h))
            rect.fill.solid()
            rect.fill.fore_color.rgb = WHITE if ridx % 2 == 0 else LIGHT
            rect.line.color.rgb = LIGHT
            tf = rect.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = cell
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(10.5)
            r.font.color.rgb = TEXT
            x += w
        cur_y += Inches(row_h)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
band.fill.solid()
band.fill.fore_color.rgb = NAVY
band.line.fill.background()
accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.05), Inches(0.18), Inches(4.7))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = ORANGE
accent_bar.line.fill.background()
add_title(slide, "Muse Reading：数据集构建方案与评测基准设计", "AI阅读系统项目汇报", dark=True)
tb = slide.shapes.add_textbox(Inches(1.15), Inches(2.0), Inches(10.8), Inches(2.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "聚焦“如何构建数据集”\n与“如何建立标准化评测基准”"
r.font.name = "Microsoft YaHei"
r.font.size = Pt(28)
r.font.bold = True
r.font.color.rgb = WHITE
tb2 = slide.shapes.add_textbox(Inches(1.18), Inches(5.6), Inches(5.0), Inches(0.35))
p2 = tb2.text_frame.paragraphs[0]
r2 = p2.add_run()
r2.text = "Muse Reading Team  |  2026.05"
r2.font.name = "Microsoft YaHei"
r2.font.size = Pt(11)
r2.font.color.rgb = LIGHT
add_footer(slide, 1)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "汇报目标与整体框架", "本次汇报只讨论数据工程与评测设计，不展开产品功能细节。")
add_bullets(
    slide,
    [
        "问题 1：项目所需数据集由哪些部分构成，如何采集、清洗、切分、标注与组织？",
        "问题 2：如何构建标准化评测样本，验证检索质量、人设一致性与防剧透安全性？",
    ],
    Inches(0.8), Inches(1.65), Inches(11.8), Inches(1.3), 17
)
add_flow_boxes(
    slide,
    ["原始数据来源", "文本清洗与解析", "语义切片与标注", "分层数据集组织", "评测集构建", "指标评估"],
    Inches(0.85), Inches(3.0), Inches(11.55), Inches(0.85), accent=TEAL
)
add_card(slide, Inches(1.0), Inches(4.35), Inches(3.6), Inches(1.7), "数据工作核心", ["强调数据“怎么建”", "而不是只讲数据“拿来做什么”"], TEAL)
add_card(slide, Inches(4.9), Inches(4.35), Inches(3.6), Inches(1.7), "工程视角", ["来源合法", "结构清晰", "元数据完整", "可复现扩展"], BLUE)
add_card(slide, Inches(8.8), Inches(4.35), Inches(3.6), Inches(1.7), "评测视角", ["样本标准化", "指标可量化", "结果可横向对比"], ORANGE)
add_footer(slide, 2)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "数据集总体设计", "构建多层次、结构化的数据资产体系。")
add_bullets(slide, ["整体上分为四类数据：基础文本语料、人物资料语料、人工标注数据、评测基准数据。"], Inches(0.82), Inches(1.4), Inches(11.6), Inches(0.6), 16)
add_card(slide, Inches(0.9), Inches(2.0), Inches(2.7), Inches(2.0), "1. 书本文本语料库", ["公版书 / 开源文本 / 合法授权 EPUB 或 TXT", "章节化存储", "语义段落切片", "附带位置与主题元数据"], BLUE)
add_card(slide, Inches(3.95), Inches(2.0), Inches(2.7), Inches(2.0), "2. 作者/角色资料库", ["作者作品", "访谈、书信、日记、传记", "开放知识源作背景补充", "按人物独立建库"], TEAL)
add_card(slide, Inches(7.0), Inches(2.0), Inches(2.7), Inches(2.0), "3. 标注型数据", ["情绪峰值", "冲突强度", "心理复杂度", "主动提示适配性"], ORANGE)
add_card(slide, Inches(10.05), Inches(2.0), Inches(2.4), Inches(2.0), "4. 评测基准", ["检索评测集", "人设一致性评测集", "防剧透评测集", "用户实验样本"], NAVY)
add_flow_boxes(slide, ["采集", "清洗", "切片", "标注", "组织", "评测"], Inches(1.15), Inches(4.75), Inches(10.9), Inches(0.72), accent=ORANGE)
add_bullets(slide, ["核心原则：统一格式、明确来源、保留上下文、强化元数据、便于后续扩展。"], Inches(0.95), Inches(5.8), Inches(11.1), Inches(0.6), 15)
add_footer(slide, 3)

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "书本文本数据集构建", "从原始书籍到可检索文本单元，重点在切片和元数据工程。")
add_flow_boxes(slide, ["原始文本获取", "文本清洗", "章节解析", "语义切片", "元数据标注", "向量化存储"], Inches(0.75), Inches(1.55), Inches(11.8), Inches(0.82), accent=BLUE)
add_card(slide, Inches(0.85), Inches(2.7), Inches(4.0), Inches(2.4), "数据来源", ["公版文学作品", "开源文本库", "合法授权 EPUB / TXT", "优先选择版权风险低、结构清晰的文本"], BLUE)
add_card(slide, Inches(5.05), Inches(2.7), Inches(3.2), Inches(2.4), "处理原则", ["去除目录噪声、页码、重复段", "按章节保留叙事边界", "以自然段/场景段切分", "避免粗暴固定 token 切块"], TEAL)
headers = ["字段", "说明"]
rows = [
    ["book_id", "书籍唯一标识"],
    ["chapter_id / section_id", "章节与小节位置"],
    ["paragraph_id", "段落编号"],
    ["characters_present", "段落涉及人物"],
    ["theme_tag / spoiler_level", "主题标签与剧透风险"],
]
add_table_like(slide, headers, rows, Inches(8.55), Inches(2.7), [1.6, 3.3], row_h=0.43)
add_bullets(slide, ["结果形态：形成可追踪、可过滤、可扩展的书本文本数据集。"], Inches(0.95), Inches(5.55), Inches(10.8), Inches(0.5), 15)
add_footer(slide, 4)

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "作者与角色资料数据集构建", "从人物材料中抽取稳定、可追溯的风格与事实信息。")
add_card(slide, Inches(0.82), Inches(1.7), Inches(3.65), Inches(2.7), "数据来源", ["作者代表作品", "公开演讲与访谈", "书信、日记、传记", "百科与开放知识源作背景补充"], NAVY)
add_card(slide, Inches(4.75), Inches(1.7), Inches(3.65), Inches(2.7), "数据组织", ["区分事实型资料与风格型资料", "保留来源标签与引用信息", "不同作者/角色独立建库", "避免人物之间语料串扰"], TEAL)
add_card(slide, Inches(8.68), Inches(1.7), Inches(3.65), Inches(2.7), "处理流程", ["原始材料收集", "去噪与去重", "分段切片", "资料标签化", "入库管理"], ORANGE)
add_flow_boxes(slide, ["事实层", "风格层", "来源层"], Inches(1.5), Inches(4.75), Inches(4.2), Inches(0.75), accent=TEAL)
add_bullets(
    slide,
    [
        "事实层：生平、时代背景、作品主题、代表观点。",
        "风格层：语言习惯、评论方式、逻辑结构、典型表达。",
        "来源层：为每条资料保留出处，保证可回溯与后续抽检。",
    ],
    Inches(6.0), Inches(4.6), Inches(6.0), Inches(1.35), 14
)
add_footer(slide, 5)

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "标注数据与主动提示数据集构建", "通过人工标注建立高质量 Golden Dataset，再扩展到全书。")
add_bullets(slide, ["目标不是随机挑选段落，而是建立一套可执行的标注规范。"], Inches(0.85), Inches(1.4), Inches(11.0), Inches(0.45), 16)
add_card(slide, Inches(0.9), Inches(2.0), Inches(2.25), Inches(2.2), "标注维度 1", ["情绪强度", "情绪转折", "叙事张力"], ORANGE)
add_card(slide, Inches(3.35), Inches(2.0), Inches(2.25), Inches(2.2), "标注维度 2", ["人际冲突", "对话密度", "关系变化"], BLUE)
add_card(slide, Inches(5.8), Inches(2.0), Inches(2.25), Inches(2.2), "标注维度 3", ["心理复杂度", "犹豫/欺骗/顿悟", "内在动机强度"], TEAL)
add_card(slide, Inches(8.25), Inches(2.0), Inches(2.25), Inches(2.2), "标注维度 4", ["象征密度", "意象出现", "主题提示价值"], NAVY)
add_card(slide, Inches(10.7), Inches(2.0), Inches(1.7), Inches(2.2), "标签", ["是否适合主动提示"], ORANGE)
add_flow_boxes(slide, ["抽取代表章节", "人工标注", "形成 Golden Dataset", "模型辅助扩展", "人工复核"], Inches(1.05), Inches(4.75), Inches(11.0), Inches(0.8), accent=ORANGE)
add_footer(slide, 6)

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "评测基准总体设计", "建立标准化测试样本，而不是只看主观观感。")
add_card(slide, Inches(0.95), Inches(1.9), Inches(2.6), Inches(2.1), "检索评测集", ["问题样本", "标准支撑段落", "query-doc 对齐关系"], BLUE)
add_card(slide, Inches(3.75), Inches(1.9), Inches(2.6), Inches(2.1), "人设一致性评测集", ["作者/角色特定问题", "标准化评分量表", "人工与模型联合判定"], TEAL)
add_card(slide, Inches(6.55), Inches(1.9), Inches(2.6), Inches(2.1), "防剧透评测集", ["前文位置问后文情节", "诱导式追问", "安全回答标注"], ORANGE)
add_card(slide, Inches(9.35), Inches(1.9), Inches(2.6), Inches(2.1), "用户实验样本", ["理解题", "回忆题", "主观问卷", "停留时长"], NAVY)
add_bullets(slide, ["评测集构建思路：先定义任务，再设计样本，再建立标注规范，最后绑定指标。"], Inches(0.95), Inches(4.55), Inches(11.4), Inches(0.55), 15)
add_flow_boxes(slide, ["任务定义", "样本设计", "人工标注", "质量抽检", "指标绑定"], Inches(1.2), Inches(5.25), Inches(10.8), Inches(0.8), accent=TEAL)
add_footer(slide, 7)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "检索与人设一致性评测基准", "同时验证“找得准”与“说得像”。")
add_card(slide, Inches(0.85), Inches(1.8), Inches(5.7), Inches(3.6), "检索评测集构建", [
    "人工设计 query，并为每个 query 标注标准支撑段落。",
    "建立 query 与目标 chapter/section/paragraph 的对应关系。",
    "分别覆盖背景解释、人物关系、主题理解等不同问题类型。",
    "指标：Top-k Recall、MRR、nDCG、命中正确章节比例。"
], BLUE)
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(3.6), "人设一致性评测集构建", [
    "针对作者/角色设计特定问答样本，覆盖评论、解释、总结等任务。",
    "建立统一评分量表：风格一致性、价值观一致性、分析框架一致性、表达自然度。",
    "评测方式：人工评分 + LLM-as-a-judge + 与普通模型输出对比。"
], ORANGE)
add_footer(slide, 8)

# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "防剧透与安全性评测基准", "通过对抗式样本验证系统是否真正做到“按阅读进度作答”。")
add_card(slide, Inches(0.85), Inches(1.8), Inches(4.1), Inches(3.5), "Spoiler-Trap Dataset 构造", [
    "在前文阅读位置提问后文关键情节。",
    "直接提问：某角色最后结局是什么？",
    "暗示提问：后面是不是会发生反转？",
    "模糊追问：他后来为什么突然改变？"
], ORANGE)
add_card(slide, Inches(5.2), Inches(1.8), Inches(3.35), Inches(3.5), "标注结果类别", [
    "安全回答",
    "错误拒答",
    "模糊泄露",
    "明确剧透"
], BLUE)
add_card(slide, Inches(8.8), Inches(1.8), Inches(3.45), Inches(3.5), "核心指标", [
    "Spoilage Leakage Rate",
    "安全回答率",
    "错误拒答率",
    "不同提问形式下的稳健性"
], TEAL)
add_bullets(slide, ["目标：把防剧透从“主观承诺”变成“可量化评测结果”。"], Inches(0.95), Inches(5.65), Inches(11.1), Inches(0.45), 15)
add_footer(slide, 9)

# Slide 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, RGBColor(242, 246, 252))
add_title(slide, "总结", "高质量数据集与标准化评测基准，决定系统上限与可信度。")
add_bullets(
    slide,
    [
        "第一，Muse Reading 需要的不是单一文本库，而是书本文本、人物资料、标注样本和评测样本构成的结构化数据体系。",
        "第二，数据集构建的关键在于：来源合法、语义切片合理、元数据完整、人物资料分层、标注规范清晰。",
        "第三，评测基准设计的关键在于：任务标准化、样本可复现、指标可量化，尤其要重点覆盖防剧透安全性。",
    ],
    Inches(0.95), Inches(1.75), Inches(11.2), Inches(2.3), 17
)
quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(4.7), Inches(11.2), Inches(1.25))
quote.fill.solid()
quote.fill.fore_color.rgb = NAVY
quote.line.fill.background()
tf = quote.text_frame
tf.clear()
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Muse Reading 的关键，不只是模型生成能力，更在于高质量数据集构建与严格评测基准设计。"
r.font.name = "Microsoft YaHei"
r.font.size = Pt(18)
r.font.bold = True
r.font.color.rgb = WHITE
add_footer(slide, 10)

out = "Muse Reading_数据集构建方案与评测基准设计.pptx"
prs.save(out)
print(out)
